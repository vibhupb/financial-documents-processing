"""Tests for compliance-ingest parser (Word, PPT, PDF)."""
import io
import os
import sys
import pytest
from unittest.mock import patch, MagicMock

sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', 'lambda', 'compliance-ingest'))
from parser import parse_docx, parse_pptx, parse_pdf, ParsedContent


def test_docx_text_and_tables():
    """parse_docx extracts paragraphs and tables."""
    from docx import Document
    doc = Document()
    doc.add_paragraph("Must specify APR.")
    t = doc.add_table(rows=2, cols=2)
    t.cell(0, 0).text = "Field"
    t.cell(0, 1).text = "Required"
    t.cell(1, 0).text = "APR"
    t.cell(1, 1).text = "Yes"
    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    result = parse_docx(buf.read())
    assert isinstance(result, ParsedContent)
    assert "APR" in result.text
    assert "Must specify APR" in result.text
    assert len(result.tables) == 1


def test_docx_images():
    """parse_docx extracts embedded images."""
    from docx import Document
    from PIL import Image
    doc = Document()
    doc.add_paragraph("Chart:")
    ib = io.BytesIO()
    Image.new('RGB', (10, 10), 'red').save(ib, format='PNG')
    ib.seek(0)
    doc.add_picture(ib, width=914400)
    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    result = parse_docx(buf.read())
    assert len(result.images) >= 1


@patch("parser.pypdf")
def test_pdf_text(mock_pypdf):
    """parse_pdf extracts text from PDF pages."""
    mock_page = MagicMock()
    mock_page.extract_text.return_value = "Page 1 content"
    mock_reader = MagicMock()
    mock_reader.pages = [mock_page]
    mock_pypdf.PdfReader.return_value = mock_reader
    result = parse_pdf(b"fake pdf bytes")
    assert "Page 1 content" in result.text
    assert "--- Page 1 ---" in result.text


def test_pptx_text():
    """parse_pptx extracts slide text."""
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Compliance"
    slide.placeholders[1].text = "APR requirement"
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    result = parse_pptx(buf.read())
    assert "Compliance" in result.text
    assert "APR" in result.text
