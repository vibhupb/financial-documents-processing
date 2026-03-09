"""Parser for Word, PPT, PDF reference documents."""
from __future__ import annotations

import io
from dataclasses import dataclass, field

import pypdf


@dataclass
class ParsedContent:
    """Parsed content from a reference document."""
    text: str = ""
    tables: list = field(default_factory=list)
    images: list = field(default_factory=list)


def parse_docx(fb: bytes) -> ParsedContent:
    """Parse a Word document, extracting text, tables, and images."""
    from docx import Document

    doc = Document(io.BytesIO(fb))
    tables = [[[c.text for c in r.cells] for r in t.rows] for t in doc.tables]
    images = [
        rel.target_part.blob
        for rel in doc.part.rels.values()
        if "image" in rel.reltype
    ]
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    for i, tbl in enumerate(tables):
        text += f"\n[Table {i + 1}]\n" + "\n".join(" | ".join(r) for r in tbl)
    return ParsedContent(text=text, tables=tables, images=images)


def parse_pptx(fb: bytes) -> ParsedContent:
    """Parse a PowerPoint presentation, extracting slide text, tables, and images."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(fb))
    parts: list[str] = []
    images: list[bytes] = []
    tables: list = []
    for n, sl in enumerate(prs.slides, 1):
        parts.append(f"--- Slide {n} ---")
        for s in sl.shapes:
            if s.has_text_frame:
                parts.append(s.text_frame.text)
            if s.has_table:
                tables.append([[c.text for c in r.cells] for r in s.table.rows])
            if s.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                images.append(s.image.blob)
    return ParsedContent(text="\n".join(parts), tables=tables, images=images)


def parse_pdf(fb: bytes) -> ParsedContent:
    """Parse a PDF document, extracting text from all pages."""
    reader = pypdf.PdfReader(io.BytesIO(fb))
    text = "\n".join(
        f"--- Page {i + 1} ---\n{p.extract_text() or ''}"
        for i, p in enumerate(reader.pages)
    )
    return ParsedContent(text=text)


def parse_xlsx(fb: bytes) -> ParsedContent:
    """Parse an Excel spreadsheet, extracting all sheet data as structured text."""
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(fb), data_only=True)
    parts: list[str] = []
    tables: list = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"--- Sheet: {sheet_name} ---")
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                rows.append(cells)
                parts.append(" | ".join(cells))
        if rows:
            tables.append(rows)
    return ParsedContent(text="\n".join(parts), tables=tables)
